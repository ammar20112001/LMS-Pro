"""
Handout ingestion pipeline.

Primary strategy: detect "Lecture N" headings in PDF to split by lecture.
Fallback strategy: when headings aren't found, use Haiku with the LMS lecture list
  to identify boundaries page-by-page with running continuity context.
"""

import json
import logging
import re
from collections import defaultdict
from pathlib import Path

import httpx

from ..db import SessionLocal
from ..models import Course, HandoutChunk, HandoutImage, HandoutSource, Lecture, utcnow
from ..config import settings

log = logging.getLogger(__name__)

HANDOUTS_DIR = Path(__file__).parents[3] / "handouts"
IMAGES_DIR = settings.data_dir / "handout_images"

_LECTURE_HEADING = re.compile(r"Lecture\s*(?:No\.?\s*)?(\d+)", re.IGNORECASE)
_SUBLECTURE = re.compile(r"^\d+\.\d+\.")   # "16.224. Title" — sub-video pattern


# ── Public entry points ───────────────────────────────────────────────────────

def run_ingestion():
    db = SessionLocal()
    try:
        _scan_and_ingest(db)
    except Exception:
        log.exception("Handout ingestion crashed")
        db.rollback()
    finally:
        db.close()


def reset_source_for_reingest(file_path: str):
    """Mark a source as pending so it will be re-ingested on the next scan."""
    db = SessionLocal()
    try:
        source = db.query(HandoutSource).filter_by(file_path=file_path).first()
        if not source:
            return
        _delete_source_chunks(db, source)
        source.ingest_status = "pending"
        source.total_chunks = 0
        db.commit()
        log.info("Reset %s for re-ingestion", Path(file_path).name)
    finally:
        db.close()


# ── Scan loop ─────────────────────────────────────────────────────────────────

def _scan_and_ingest(db):
    if not HANDOUTS_DIR.exists():
        log.warning("handouts/ directory not found at %s", HANDOUTS_DIR)
        return

    files = sorted(
        list(HANDOUTS_DIR.glob("**/*.pdf")) + list(HANDOUTS_DIR.glob("**/*.pptx"))
    )
    log.info("Ingestion scan: found %d handout files", len(files))

    for file_path in files:
        course_code = _course_code_from_path(file_path)
        if not course_code:
            continue

        source = db.query(HandoutSource).filter_by(file_path=str(file_path)).first()
        if source and source.ingest_status == "done":
            continue

        if not source:
            source = HandoutSource(
                course_code=course_code,
                file_path=str(file_path),
                file_type=file_path.suffix.lstrip(".").lower(),
            )
            db.add(source)
            db.flush()

        log.info("Ingesting %s (%s)", file_path.name, course_code)
        source.ingest_status = "ingesting"
        db.commit()

        try:
            if source.file_type == "pdf":
                chunks = _ingest_pdf(db, source, file_path, course_code)
            else:
                chunks = _ingest_pptx(db, source, file_path, course_code)

            source.total_chunks = chunks
            source.ingest_status = "done"
            source.ingested_at = utcnow()
            db.commit()
            log.info("Ingested %s — %d chunks", file_path.name, chunks)
        except Exception as e:
            log.exception("Failed to ingest %s: %s", file_path.name, e)
            source.ingest_status = "failed"
            db.commit()


# ── PDF ingestion ─────────────────────────────────────────────────────────────

def _ingest_pdf(db, source: HandoutSource, file_path: Path, course_code: str) -> int:
    import fitz

    doc = fitz.open(str(file_path))
    pages_text = [page.get_text("text") for page in doc]
    total_pages = len(pages_text)

    boundaries = _detect_lecture_boundaries(pages_text)

    if not boundaries:
        # Primary strategy failed — try Haiku fallback
        lms_lectures = _get_lms_lectures(db, course_code)
        if lms_lectures:
            log.info("%s: no headings found, running Haiku fallback (%d LMS lectures)",
                     file_path.name, len(lms_lectures))
            doc.close()
            return _ingest_pdf_fallback(db, source, file_path, course_code, lms_lectures)

        # No LMS lectures either — single chunk
        raw_text = "\n".join(pages_text)
        chunk = _save_chunk(db, source, course_code, 1, file_path.stem, raw_text, 0, total_pages - 1)
        _extract_pdf_images(db, doc, 0, total_pages - 1, chunk.id, course_code)
        doc.close()
        return 1

    chunk_count = 0
    for idx, (page_start, lec_no, title) in enumerate(boundaries):
        page_end = boundaries[idx + 1][0] - 1 if idx + 1 < len(boundaries) else total_pages - 1
        raw_text = "\n".join(pages_text[page_start:page_end + 1])
        chunk = _save_chunk(db, source, course_code, lec_no, title, raw_text, page_start, page_end)
        _extract_pdf_images(db, doc, page_start, page_end, chunk.id, course_code)
        chunk_count += 1

    doc.close()
    return chunk_count


def _detect_lecture_boundaries(pages_text: list[str]) -> list[tuple[int, int, str]]:
    """Return list of (page_idx, lecture_no, title), deduplicated by lecture number."""
    raw: list[tuple[int, int, str]] = []
    for i, text in enumerate(pages_text):
        m = _LECTURE_HEADING.search(text)
        if m:
            lec_no = int(m.group(1))
            raw.append((i, lec_no, _extract_title_from_page(text, lec_no)))

    seen: set[int] = set()
    result: list[tuple[int, int, str]] = []
    for b in raw:
        if b[1] not in seen:
            seen.add(b[1])
            result.append(b)
    return result


# ── Haiku fallback ────────────────────────────────────────────────────────────

def _ingest_pdf_fallback(db, source: HandoutSource, file_path: Path,
                         course_code: str, lms_lectures: list[tuple[int, str]]) -> int:
    """
    When the primary heading-detection finds nothing, send pages in batches to Haiku.
    Each batch gets the running summary of the previous batch as continuity context.
    Haiku returns where new lectures start; we accumulate pages per lecture and build chunks.
    """
    import fitz

    doc = fitz.open(str(file_path))
    pages_text = [doc[i].get_text("text") for i in range(len(doc))]
    total_pages = len(pages_text)

    # page_idx -> lecture_no mapping
    page_lecture: dict[int, int] = {}
    current_lec_no = lms_lectures[0][0]
    prev_summary = ""

    BATCH = 8
    for batch_start in range(0, total_pages, BATCH):
        batch = pages_text[batch_start: batch_start + BATCH]
        transitions, prev_summary = _haiku_detect_boundaries(
            batch, batch_start, lms_lectures, current_lec_no, prev_summary
        )

        # Apply transitions page by page within this batch
        t_map = {t["at_page"]: t["lecture_no"] for t in transitions}
        for offset in range(len(batch)):
            if offset + 1 in t_map:
                current_lec_no = t_map[offset + 1]
            page_lecture[batch_start + offset] = current_lec_no

        log.debug("Batch %d-%d: %d transitions, now lec %d",
                  batch_start + 1, batch_start + len(batch), len(transitions), current_lec_no)

    # Group page indices by lecture
    grouped: dict[int, list[int]] = defaultdict(list)
    for page_idx in range(total_pages):
        grouped[page_lecture.get(page_idx, current_lec_no)].append(page_idx)

    # Build a quick lookup for LMS lecture titles
    title_map = dict(lms_lectures)

    # Delete any previously created chunks for this source before rebuilding
    _delete_source_chunks(db, source)

    chunk_count = 0
    for lec_no in sorted(grouped.keys()):
        page_indices = grouped[lec_no]
        raw_text = "\n".join(pages_text[i] for i in page_indices)
        lms_title = title_map.get(lec_no, "")
        title = f"Lecture {lec_no}" + (f" — {lms_title}" if lms_title else "")
        chunk = _save_chunk(
            db, source, course_code, lec_no, title, raw_text,
            min(page_indices), max(page_indices)
        )
        _extract_pdf_images(db, doc, min(page_indices), max(page_indices), chunk.id, course_code)
        chunk_count += 1

    doc.close()
    db.flush()
    log.info("Fallback ingestion: %s → %d lecture chunks", Path(source.file_path).name, chunk_count)
    return chunk_count


def _haiku_detect_boundaries(
    batch_pages: list[str],
    batch_start: int,
    lms_lectures: list[tuple[int, str]],
    current_lec_no: int,
    prev_summary: str,
) -> tuple[list[dict], str]:
    """
    Send one batch of pages to Haiku. Returns (transitions, summary_for_next_batch).
    transitions: [{"at_page": <1-based within batch>, "lecture_no": <int>}]
    """
    lecture_list_str = "\n".join(
        f"{no}. {title}" for no, title in lms_lectures[:60]
    )
    pages_str = "\n---PAGE BREAK---\n".join(
        f"[Page {batch_start + i + 1}]\n{text[:900]}"
        for i, text in enumerate(batch_pages)
    )

    prompt = (
        f"Course lecture list:\n{lecture_list_str}\n\n"
        f"Currently on lecture: {current_lec_no}\n"
        f"Previous content summary: {prev_summary or 'Start of document'}\n\n"
        f"Pages {batch_start + 1}–{batch_start + len(batch_pages)}:\n{pages_str[:7000]}\n\n"
        "Identify where new lectures begin in this batch of pages.\n"
        "Return ONLY valid JSON:\n"
        '{"transitions":[{"at_page":<1-based position in batch>,"lecture_no":<int>}],'
        '"summary":"<50-word summary of these pages for continuity>"}\n'
        f"If all pages continue lecture {current_lec_no}, return empty transitions array."
    )

    try:
        resp = httpx.post(
            "https://api.anthropic.com/v1/messages",
            headers={
                "x-api-key": settings.anthropic_api_key,
                "anthropic-version": "2023-06-01",
                "content-type": "application/json",
            },
            json={
                "model": "claude-haiku-4-5-20251001",
                "max_tokens": 400,
                "system": "You are a document structure analyzer. Return only valid JSON, no other text.",
                "messages": [{"role": "user", "content": prompt}],
            },
            timeout=60,
        )
        resp.raise_for_status()
        raw = resp.json()["content"][0]["text"]
        data = json.loads(raw[raw.index("{"):raw.rindex("}") + 1])
        return data.get("transitions", []), data.get("summary", "")
    except Exception as e:
        log.warning("Haiku boundary detection failed for batch %d: %s", batch_start, e)
        return [], ""


# ── LMS lecture list helper ───────────────────────────────────────────────────

def _get_lms_lectures(db, course_code: str) -> list[tuple[int, str]]:
    """
    Return (serial_no, title) pairs from the lectures table for this course.
    Deduplicates by serial_no (keeps first) and drops sub-lecture entries like '16.224. Title'.
    """
    course = db.query(Course).filter_by(code=course_code).first()
    if not course:
        return []

    seen: set[int] = set()
    result: list[tuple[int, str]] = []
    for lec in sorted(course.lectures, key=lambda l: l.serial_no):
        if lec.serial_no in seen:
            continue
        if _SUBLECTURE.match(lec.title or ""):
            continue   # skip "16.224. Designing Polyglot Services" style sub-videos
        seen.add(lec.serial_no)
        result.append((lec.serial_no, lec.title or ""))
    return result


# ── PPTX ingestion (unchanged) ────────────────────────────────────────────────

def _ingest_pptx(db, source: HandoutSource, file_path: Path, course_code: str) -> int:
    from pptx import Presentation

    lec_range = _lecture_range_from_filename(file_path.stem)
    lec_start, lec_end = lec_range if lec_range else (1, 1)

    prs = Presentation(str(file_path))
    slides = list(prs.slides)
    total_slides = len(slides)
    if not slides:
        return 0

    lec_count = lec_end - lec_start + 1
    slides_per_lec = max(1, total_slides // lec_count)
    out_dir = IMAGES_DIR / course_code
    out_dir.mkdir(parents=True, exist_ok=True)

    chunk_count = 0
    for lec_offset in range(lec_count):
        lec_no = lec_start + lec_offset
        slide_start = lec_offset * slides_per_lec
        slide_end = slide_start + slides_per_lec if lec_offset < lec_count - 1 else total_slides
        lec_slides = slides[slide_start:slide_end]

        text_parts = []
        for slide in lec_slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(r.text for r in para.runs).strip()
                        if line:
                            text_parts.append(line)

        raw_text = "\n".join(text_parts)
        chunk = _save_chunk(db, source, course_code, lec_no, f"Lecture {lec_no}",
                            raw_text, slide_start, slide_end - 1)

        seq = 1
        for slide in lec_slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:
                    try:
                        img = shape.image
                        out_path = out_dir / f"{chunk.id}_{seq}.{img.ext}"
                        out_path.write_bytes(img.blob)
                        db.add(HandoutImage(chunk_id=chunk.id, seq=seq, file_path=str(out_path)))
                        seq += 1
                    except Exception as e:
                        log.debug("PPTX image extract error: %s", e)
        if seq > 1:
            db.flush()
        chunk_count += 1

    return chunk_count


# ── Shared utilities ──────────────────────────────────────────────────────────

def _course_code_from_path(file_path: Path) -> str | None:
    for candidate in (file_path.stem.upper(), file_path.parent.name.upper()):
        m = re.match(r"(CS\d{3,4})", candidate)
        if m:
            return m.group(1)
    return None


def _extract_title_from_page(text: str, lec_no: int) -> str:
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    for i, line in enumerate(lines[:15]):
        if re.search(rf"lecture\s*(?:no\.?\s*)?{lec_no}", line, re.IGNORECASE):
            if re.search(r"\.{3,}|\s\d{1,3}\s*$", line):
                continue
            next_lines = [
                l for l in lines[i + 1:i + 4]
                if len(l) > 3 and not re.search(r"\.{3,}|\s\d{1,3}\s*$", l)
            ]
            if next_lines:
                return f"Lecture {lec_no} — {next_lines[0]}"
            return f"Lecture {lec_no}"
    return f"Lecture {lec_no}"


def _extract_pdf_images(db, doc, page_start: int, page_end: int, chunk_id: int, course_code: str):
    out_dir = IMAGES_DIR / course_code
    out_dir.mkdir(parents=True, exist_ok=True)
    seq = 1
    for page_idx in range(page_start, page_end + 1):
        page = doc[page_idx]
        for img_info in page.get_images(full=True):
            xref = img_info[0]
            try:
                img_bytes = doc.extract_image(xref)
                if img_bytes["width"] < 50 or img_bytes["height"] < 50:
                    continue
                ext = img_bytes["ext"]
                out_path = out_dir / f"{chunk_id}_{seq}.{ext}"
                out_path.write_bytes(img_bytes["image"])
                db.add(HandoutImage(chunk_id=chunk_id, seq=seq, file_path=str(out_path)))
                seq += 1
            except Exception as e:
                log.debug("Image extraction error xref=%d: %s", xref, e)
    if seq > 1:
        db.flush()


def _delete_source_chunks(db, source: HandoutSource):
    """Delete all chunks and their images for a source (disk + DB)."""
    chunks = db.query(HandoutChunk).filter_by(source_id=source.id).all()
    for chunk in chunks:
        for img in chunk.images:
            try:
                Path(img.file_path).unlink(missing_ok=True)
            except Exception:
                pass
        db.delete(chunk)
    db.flush()


def _lecture_range_from_filename(stem: str) -> tuple[int, int] | None:
    m = re.search(r"(\d+)[_\-](\d+)", stem)
    if m:
        return int(m.group(1)), int(m.group(2))
    m = re.search(r"(\d+)", stem)
    if m:
        n = int(m.group(1))
        return n, n
    return None


def _save_chunk(db, source: HandoutSource, course_code: str, lecture_no: int,
                title: str, raw_text: str, page_start: int, page_end: int) -> HandoutChunk:
    chunk = HandoutChunk(
        source_id=source.id,
        course_code=course_code,
        lecture_no=lecture_no,
        title=title,
        raw_text=raw_text,
        page_start=page_start,
        page_end=page_end,
        enrich_status="pending",
    )
    db.add(chunk)
    db.flush()
    return chunk
